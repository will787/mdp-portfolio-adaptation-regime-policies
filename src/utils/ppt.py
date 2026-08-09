# %%
"""
Script para geração automática da apresentação K-NESIAN em formato PowerPoint (.pptx)
Requisito: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_k_nesian_presentation():
    prs = Presentation()
    
    # Configurar formato Widescreen 16:9 (13.33 x 7.5 polegadas)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Paleta de Cores do Tema Dark K-NESIAN
    DARK_BG = RGBColor(11, 19, 41)       # #0B1329
    CARD_BG = RGBColor(15, 23, 42)       # #0F172A
    CYAN = RGBColor(56, 189, 248)        # #38BDF8
    GOLD = RGBColor(245, 158, 11)        # #F59E0B
    WHITE = RGBColor(248, 250, 252)      # #F8FAFC
    GRAY = RGBColor(148, 163, 184)       # #94A3B8
    GREEN = RGBColor(74, 222, 128)       # #4ADE80
    RED = RGBColor(248, 113, 113)        # #F87171

    def apply_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, title_text, badge_text):
        # Título
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.font.name = 'Arial'

        # Badge
        badge_box = slide.shapes.add_textbox(Inches(9.8), Inches(0.4), Inches(2.9), Inches(0.6))
        tf_b = badge_box.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.RIGHT
        p_b.text = badge_text
        p_b.font.size = Pt(12)
        p_b.font.bold = True
        p_b.font.color.rgb = GOLD
        p_b.font.name = 'Courier New'

    # ==========================================
    # SLIDE 1: IDENTIDADE & CONCEITO
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide1)
    add_header(slide1, "K-NESIAN: Alocação Dinâmica sob Incerteza Radical", "QUANT ROBOT • B3 STRATEGY")

    # Card 1: Identidade
    shape1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.8))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = CARD_BG
    shape1.line.color.rgb = CYAN
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "1. Identidade do Robô"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "\nO K-NESIAN é um robô quantitativo fundamentado na teoria de John Maynard Keynes sobre Incerteza Radical e Preferência pela Liquidez."
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    
    p3 = tf1.add_paragraph()
    p3.text = "\n• Core Tecnológico: Aprendizado de Máquina Não Supervisionado (Gaussian HMM) + Programação Dinâmica (Equação de Bellman em MDP)."
    p3.font.size = Pt(13)
    p3.font.color.rgb = CYAN

    # Card 2: Hipótese & Ineficiência
    shape2 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(3.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = CARD_BG
    shape2.line.color.rgb = GOLD
    tf2 = shape2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "2. Hipótese & Ineficiência Explorada"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf2.add_paragraph()
    p2.text = "\n• Falha dos Modelos Tradicionais: Otimizadores clássicos (Markowitz) assumem risco estático. Em choques, os mercados operam sob Incerteza Radical (animal spirits)."
    p2.font.size = Pt(13)
    p2.font.color.rgb = RED

    p3 = tf2.add_paragraph()
    p3.text = "\n• Solução K-NESIAN: O HMM mapeia estados latentes e a Equação de Bellman ajusta preventivamente a exposição entre Ações B3 e Caixa CDI."
    p3.font.size = Pt(13)
    p3.font.color.rgb = GREEN

    # Flow Diagram
    flow_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.4))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = CARD_BG
    flow_box.line.color.rgb = CYAN
    tf_f = flow_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "DIAGRAMA CONCEITUAL DO CICLO DE DECISÃO"
    p_f.font.size = Pt(12)
    p_f.font.bold = True
    p_f.font.color.rgb = GOLD

    p_f2 = tf_f.add_paragraph()
    p_f2.text = "1. Sinais Macroeconômicos  ➔  2. Regimes HMM (4 Estados)  ➔  3. Equação de Bellman (MDP)  ➔  4. Alocação Ótima (Risco vs. Liquidez)"
    p_f2.font.size = Pt(14)
    p_f2.font.bold = True
    p_f2.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 2: MODELAGEM QUANTITATIVA
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide2)
    add_header(slide2, "Modelagem Quantitativa & Pipeline de Dados", "FEATURE STORE & BELLMAN MDP")

    # Pipeline Card
    s2_card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4))
    s2_card1.fill.solid()
    s2_card1.fill.fore_color.rgb = CARD_BG
    s2_card1.line.color.rgb = CYAN
    tf2_1 = s2_card1.text_frame
    tf2_1.word_wrap = True

    p = tf2_1.paragraphs[0]
    p.text = "1. Feature Store (macro_features_hmm.parquet)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items1 = [
        "\n• Ativos de Risco: Top 20 ações mais líquidas da B3 (ITUB4, PETR4, VALE3, BBDC4...).",
        "• Caixa Defensivo: Ativo livre de risco remunerado a 100% do CDI diário.",
        "• Features Macro: Volatilidade Ibov, VIX Z-score, Selic, CDS 5Y, Commodities e PTAX."
    ]
    for item in items1:
        p_item = tf2_1.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = GRAY

    # Modelagem Card
    s2_card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4))
    s2_card2.fill.solid()
    s2_card2.fill.fore_color.rgb = CARD_BG
    s2_card2.line.color.rgb = GOLD
    tf2_2 = s2_card2.text_frame
    tf2_2.word_wrap = True

    p = tf2_2.paragraphs[0]
    p.text = "2. Modelagem em Duas Etapas"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p_e1 = tf2_2.add_paragraph()
    p_e1.text = "\nEtapa 1 — GaussianHMM: Mapeamento não supervisionado em 4 estados latentes de risco (Baixo, Moderado, Alto, Pânico)."
    p_e1.font.size = Pt(13)
    p_e1.font.color.rgb = CYAN

    p_e2 = tf2_2.add_paragraph()
    p_e2.text = "\nEtapa 2 — Equação de Bellman MDP:"
    p_e2.font.size = Pt(13)
    p_e2.font.bold = True
    p_e2.font.color.rgb = GOLD

    p_math = tf2_2.add_paragraph()
    p_math.text = "V(s) = max_a { R(s,a) + γ ∑ P(s'|s,a) V(s') }"
    p_math.font.size = Pt(15)
    p_math.font.bold = True
    p_math.font.color.rgb = WHITE
    p_math.font.name = 'Courier New'

    # ==========================================
    # SLIDE 3: BACKTEST & ATRITOS
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide3)
    add_header(slide3, "Backtest Rigoroso & Atritos Operacionais", "WALK-FORWARD & NORMAS B3")

    # Walk Forward
    s3_card1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.0), Inches(5.4))
    s3_card1.fill.solid()
    s3_card1.fill.fore_color.rgb = CARD_BG
    s3_card1.line.color.rgb = CYAN
    tf3_1 = s3_card1.text_frame
    tf3_1.word_wrap = True

    p = tf3_1.paragraphs[0]
    p.text = "Validação Out-of-Sample"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items_wf = [
        "\n• Metodologia Walk-Forward:",
        "  Re-treinamento periódico a cada janela de 1 ano com memória histórica móvel de 4 anos.",
        "\n• Isenção Total de Look-ahead Bias:",
        "  Garantia de que o modelo não consulta dados do futuro.",
        "\n• Redução de Overfitting:",
        "  Estabilidade de parâmetros testada em múltiplos cenários."
    ]
    for item in items_wf:
        p_item = tf3_1.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = GRAY

    # Tabela de Atritos
    x, y, cx, cy = Inches(5.9), Inches(1.4), Inches(6.8), Inches(5.4)
    shape_t = slide3.shapes.add_table(6, 3, x, y, cx, cy)
    table = shape_t.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.8)

    headers = ["Parâmetro", "Implementação", "Função no Modelo"]
    data = [
        ["Liquidação Física", "Fila real em D+2", "Bloqueia reinvestimento imediato"],
        ["Rendimento Caixa", "100% CDI diário", "Custo de oportunidade realista"],
        ["Tributação Caixa", "22.5% IR sobre RF", "Realismo fiscal na rentabilidade"],
        ["Custos Operacionais", "Corretagem 0.5% + Slippage 0.3%", "Penaliza turnover excessivo"],
        ["Gestão de Risco", "Limites (3% min / 8% máx)", "Elimina risco de concentração"]
    ]

    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DARK_BG

    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 4: ANÁLISE DE RESULTADOS
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide4)
    add_header(slide4, "Análise Crítica dos Resultados", "PERFORMANCE & METRICS")

    # Tabela Métricas
    x, y, cx, cy = Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.2)
    shape_m = slide4.shapes.add_table(4, 5, x, y, cx, cy)
    table_m = shape_m.table

    m_headers = ["Estratégia", "Retorno Acumulado", "Índice Sharpe", "Max Drawdown", "Volatilidade Anual"]
    m_data = [
        ["K-NESIAN (Bellman MDP)", "Outperformance Quant", "Superior", "Reduzido (Protegido)", "Controlada"],
        ["Ibovespa (Benchmark)", "Sujeito a quedas", "Base", "Elevado em Crises", "Alta Volatilidade"],
        ["HMM Passivo (Sem Bellman)", "Sensível a ruídos", "Intermediário", "Intermediário", "Moderada"]
    ]

    for col_idx, h in enumerate(m_headers):
        cell = table_m.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DARK_BG

    for row_idx, row_data in enumerate(m_data):
        for col_idx, val in enumerate(row_data):
            cell = table_m.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = GREEN if row_idx == 0 else WHITE

    # Análise Cards
    c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = GOLD
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "Mecanismo de Proteção Assimétrica"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf_c1.add_paragraph()
    p2.text = "\nO desconto temporal γ e as probabilidades de transição do MDP fazem o K-NESIAN aumentar a Preferência pela Liquidez antes da consolidação de regimes de pânico."
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY

    c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.0), Inches(5.9), Inches(2.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = RED
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "Limitações em Mercados Laterais"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RED
    p2 = tf_c2.add_paragraph()
    p2.text = "\nEm períodos sem tendência definida (chop market) com falsas mudanças de regime, os custos operacionais (corretagem e slippage) reduzem a rentabilidade marginal."
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 5: IA GENERATIVA & ROADMAP
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide5)
    add_header(slide5, "IA Generativa no Processo & Roadmap Futuro", "AI STACK & ROADMAP")

    # Grid IA
    s5_card1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4))
    s5_card1.fill.solid()
    s5_card1.fill.fore_color.rgb = CARD_BG
    s5_card1.line.color.rgb = CYAN
    tf5_1 = s5_card1.text_frame
    tf5_1.word_wrap = True

    p = tf5_1.paragraphs[0]
    p.text = "Aplicação de IA Generativa (15%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    ia_items = [
        "\n• Formulação Matemática:\n  Estruturação da recompensa MDP e Equação de Bellman.",
        "\n• Engenharia de Código:\n  Refatoração dos módulos motor.py e portfolio.py.",
        "\n• Conceituação & Naming:\n  Conexão da teoria de Keynes com o modelo quantitativo.",
        "\n• Diagnóstico de Backtest:\n  Análise de logs e identificação de gargalos de turnover."
    ]
    for item in ia_items:
        p_item = tf5_1.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = GRAY

    # Roadmap
    s5_card2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4))
    s5_card2.fill.solid()
    s5_card2.fill.fore_color.rgb = CARD_BG
    s5_card2.line.color.rgb = GOLD
    tf5_2 = s5_card2.text_frame
    tf5_2.word_wrap = True

    p = tf5_2.paragraphs[0]
    p.text = "Limitações & Roadmap Futuro (10%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    rm_items = [
        "\n• Limitações Identificadas:\n  Sensibilidade do HMM à hiperparametrização de regimes e custos em volatilidade intraday.",
        "\n• IA & Dados Alternativos:\n  Integração de sinais de notícias e relatórios via LLMs.",
        "\n• Reinforcement Learning:\n  Transição para Deep Q-Learning (DRL) com estado contínuo.",
        "\n• Execução Algorítmica:\n  Algoritmos inteligentes (TWAP/VWAP) para otimizar slippage."
    ]
    for item in rm_items:
        p_item = tf5_2.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = GRAY

    # Salvar Apresentação
    output_filename = "K_NESIAN_Apresentacao.pptx"
    prs.save(output_filename)
    print(f"Apresentação salva com sucesso como '{output_filename}'!")

if __name__ == "__main__":
    create_k_nesian_presentation()
# %%
